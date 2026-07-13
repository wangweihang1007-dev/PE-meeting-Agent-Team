from __future__ import annotations

import copy
import json
import re
from datetime import datetime
from pathlib import Path
from typing import Any

from openpyxl import load_workbook

from deepseek_client import DeepSeekClient


FIELD_ORDER = [
    "序号",
    "项目简称",
    "成立时间",
    "城市",
    "上市申报预期",
    "前一轮融资时点和估值",
    "已投机构",
    "本轮投前估值",
    "本次融资额（或投后估值）",
    "本次融资截止时间",
    "主营业务",
    "价值",
    "收入",
    "利润",
    "项目来源和录入时间",
    "备注",
    "否决原因",
]

REQUIRED_SECTIONS = ["1.团队", "2.股权结构", "3.产品", "4.技术", "5.生产、客户", "6.市场", "7.收入"]


def _read_text_file(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    raise ValueError(f"无法读取文本文件：{path}")


def extract_material_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        return _read_text_file(path)
    if suffix == ".docx":
        from docx import Document

        document = Document(str(path))
        parts: list[str] = []
        parts.extend(p.text for p in document.paragraphs if p.text.strip())
        for table in document.tables:
            for row in table.rows:
                values = [cell.text.strip() for cell in row.cells]
                if any(values):
                    parts.append(" | ".join(values))
        return "\n".join(parts)
    if suffix == ".pdf":
        import pdfplumber

        pages: list[str] = []
        with pdfplumber.open(str(path)) as pdf:
            for index, page in enumerate(pdf.pages, 1):
                text = page.extract_text(layout=True) or page.extract_text() or ""
                if text.strip():
                    pages.append(f"## 第{index}页\n{text.strip()}")
        return "\n\n".join(pages)
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(str(path))
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        values = [cell.text.strip() for cell in row.cells]
                        if any(values):
                            texts.append(" | ".join(values))
            if texts:
                slides.append(f"## Slide {index}\n" + "\n".join(texts))
        return "\n\n".join(slides)
    if suffix in {".xlsx", ".xlsm"}:
        workbook = load_workbook(path, data_only=True)
        sheets: list[str] = []
        for ws in workbook.worksheets:
            rows: list[str] = []
            for row in ws.iter_rows():
                values = ["" if cell.value is None else str(cell.value) for cell in row]
                while values and not values[-1]:
                    values.pop()
                if any(value.strip() for value in values):
                    rows.append(" | ".join(values))
            if rows:
                sheets.append(f"## {ws.title}\n" + "\n".join(rows))
        return "\n\n".join(sheets)
    raise ValueError(f"暂不支持的项目录入材料类型：{path}")


def _extract_json_object(text: str) -> dict[str, Any]:
    cleaned = text.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?", "", cleaned).strip()
        cleaned = re.sub(r"```$", "", cleaned).strip()
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", cleaned, re.S)
        if not match:
            raise
        return json.loads(match.group(0))


def build_project_intake_prompt(
    *,
    final_summary: str,
    background: str,
    qcc_materials: list[tuple[str, str]],
    other_materials: list[tuple[str, str]],
    project_source: str,
    intake_date: str,
) -> str:
    qcc_block = "\n\n".join(f"# {name}\n{text}" for name, text in qcc_materials) or "未提供"
    other_block = "\n\n".join(f"# {name}\n{text}" for name, text in other_materials) or "未提供"
    return f"""
你是冯源资本投资项目表录入智能体。请根据资料生成新版项目录入 JSON。

【严格边界】
1. 企查查/工商资料只用于项目录入，不用于改写会议纪要。
2. 成立时间、城市、股权、实控人优先使用企查查/工商资料。
3. 经营、融资、财务、客户和产品现状优先使用最新会议纪要；BP用于补充产品、技术、历史和市场。
4. 不臆测。没有可靠信息的字段留空。
5. 不把发货额写成收入，不把意向客户写成客户，不把送样写成量产，不把规划产能写成现有产能。
6. 冲突必须写入“备注”。

【新版字段】
只输出一个 JSON 对象，字段如下：
项目简称、成立时间、城市、上市申报预期、前一轮融资时点和估值、已投机构、本轮投前估值、本次融资额（或投后估值）、本次融资截止时间、主营业务、价值、收入、利润、项目来源和录入时间、备注、否决原因。
不要输出“序号”，不要输出“是否通过”。

【价值栏格式】
“价值”字段必须按以下七段，顺序不能变：
1.团队
2.股权结构
3.产品
4.技术
5.生产、客户
6.市场
7.收入

【写法】
- 主营业务写具体产品/解决方案，不写宽泛赛道。
- 使用白描和低 bit 写法，优先数字、年份、状态。
- 未来事项加“预计/计划/目标/公司表示”。
- “备注”必须写资料来源组合，例如：资料来源：BP、纪要、QCC。若有冲突，追加“冲突：...”。
- “项目来源和录入时间”写：项目来源：{project_source}\\n录入时间：{intake_date}

【最终会议纪要】
{final_summary}

【BP/公司材料】
{background}

【企查查/工商资料，仅用于项目录入】
{qcc_block}

【其他补充资料】
{other_block}
""".strip()


def validate_intake_entry(entry: dict[str, Any]) -> list[str]:
    warnings: list[str] = []
    value = str(entry.get("价值", ""))
    positions = [value.find(section) for section in REQUIRED_SECTIONS]
    missing = [section for section, pos in zip(REQUIRED_SECTIONS, positions) if pos < 0]
    if missing:
        warnings.append("价值栏缺少章节：" + "、".join(missing))
    elif positions != sorted(positions):
        warnings.append("价值栏章节顺序异常")
    if "是否通过" in entry:
        warnings.append("新版表格已取消“是否通过”，已忽略该字段")
    if not str(entry.get("备注", "")).strip():
        warnings.append("备注为空，应标注资料来源")
    return warnings


def generate_project_intake(
    *,
    client: DeepSeekClient,
    final_summary: str,
    background: str,
    qcc_paths: list[Path],
    other_paths: list[Path],
    project_source: str,
    intake_date: str,
) -> tuple[dict[str, Any], list[str]]:
    qcc_materials = [(path.name, extract_material_text(path)) for path in qcc_paths]
    other_materials = [(path.name, extract_material_text(path)) for path in other_paths]
    prompt = build_project_intake_prompt(
        final_summary=final_summary,
        background=background,
        qcc_materials=qcc_materials,
        other_materials=other_materials,
        project_source=project_source,
        intake_date=intake_date,
    )
    result = client.chat(
        [
            {"role": "system", "content": "你只输出合法 JSON，不输出解释。"},
            {"role": "user", "content": prompt},
        ],
        temperature=0.1,
    )
    entry = _extract_json_object(result)
    warnings = validate_intake_entry(entry)
    return entry, warnings


def _parse_date(value: Any) -> Any:
    if not isinstance(value, str):
        return value
    text = value.strip()
    for fmt in ("%Y-%m-%d", "%Y/%m/%d", "%Y.%m.%d"):
        try:
            return datetime.strptime(text, fmt)
        except ValueError:
            pass
    return value


def _is_formal_row(ws, row: int) -> bool:
    seq = ws.cell(row, 1).value
    name = ws.cell(row, 2).value
    if not name:
        return False
    return isinstance(seq, int) or (isinstance(seq, str) and seq.strip().isdigit())


def _last_formal_row(ws) -> int:
    last = 0
    for row in range(1, ws.max_row + 1):
        if _is_formal_row(ws, row):
            last = row
    if not last:
        raise ValueError("未找到可复制格式的正式项目行")
    return last


def _next_sequence(ws, last_row: int) -> int:
    values: list[int] = []
    for row in range(1, last_row + 1):
        seq = ws.cell(row, 1).value
        if isinstance(seq, int):
            values.append(seq)
        elif isinstance(seq, str) and seq.strip().isdigit():
            values.append(int(seq.strip()))
    return (max(values) if values else 0) + 1


def write_intake_excel(entry: dict[str, Any], template_path: Path, output_path: Path) -> None:
    workbook = load_workbook(template_path)
    ws = workbook.worksheets[0]
    source_row = _last_formal_row(ws)
    target_row = source_row + 1
    ws.row_dimensions[target_row].height = ws.row_dimensions[source_row].height
    for col in range(1, len(FIELD_ORDER) + 1):
        src = ws.cell(source_row, col)
        dst = ws.cell(target_row, col)
        if src.has_style:
            dst._style = copy.copy(src._style)
        dst.number_format = src.number_format
        dst.alignment = copy.copy(src.alignment)
        dst.protection = copy.copy(src.protection)

    for col, field in enumerate(FIELD_ORDER, 1):
        if field == "序号":
            value = _next_sequence(ws, source_row)
        elif field == "成立时间":
            value = _parse_date(entry.get(field, ""))
        else:
            value = entry.get(field, "")
        ws.cell(target_row, col).value = value

    output_path.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(output_path)
